import io
import logging
import re
from pathlib import Path

from pydantic import BaseModel

logger = logging.getLogger(__name__)


class ExtractedDocument(BaseModel):
    """Result of content extraction from a file."""

    title: str
    text: str
    sections: list[dict] = []  # Optional, for future use


def extract_text_plain(path: str) -> ExtractedDocument:
    """Extract content from a plain text file."""
    encodings = ["utf-8", "latin-1", "cp1252"]

    for encoding in encodings:
        try:
            with open(path, "r", encoding=encoding) as f:
                text = f.read()
            title = Path(path).stem
            return ExtractedDocument(title=title, text=text)
        except UnicodeDecodeError:
            continue

    raise ValueError(f"Could not decode file with any supported encoding: {path}")


def extract_pdf(path: str) -> ExtractedDocument:
    """Extract content from a PDF file using pdfminer.six."""
    from pdfminer.high_level import extract_text

    text = extract_text(path)
    title = Path(path).stem
    return ExtractedDocument(title=title, text=text)


def extract_docx(path: str) -> ExtractedDocument:
    """Extract content from a DOCX file using python-docx."""
    from docx import Document

    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs]
    text = "\n".join(paragraphs)

    # Try to get title from document properties
    title = doc.core_properties.title or Path(path).stem

    return ExtractedDocument(title=title, text=text)


def extract_pptx(path: str) -> ExtractedDocument:
    """Extract content from a PPTX file using python-pptx."""
    from pptx import Presentation

    prs = Presentation(path)
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_content.append(shape.text)
        if slide_content:
            slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_content))

    text = "\n\n".join(slides_text)
    title = Path(path).stem

    return ExtractedDocument(title=title, text=text)


# MIME type to extractor mapping
EXTRACTORS = {
    "text/plain": extract_text_plain,
    "application/pdf": extract_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_pptx,
}

# File extension fallbacks
EXTENSION_MIME_MAP = {
    ".txt": "text/plain",
    ".md": "text/plain",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def extract_content(path: str, file_type: str) -> ExtractedDocument:
    """
    Extract text content from a file based on its type.

    Args:
        path: Path to the file
        file_type: MIME type of the file

    Returns:
        ExtractedDocument with title and text

    Raises:
        ValueError: If file type is not supported
        Exception: If extraction fails
    """
    # Try to find extractor by MIME type
    extractor = EXTRACTORS.get(file_type)

    # Fall back to extension-based detection
    if not extractor:
        ext = Path(path).suffix.lower()
        fallback_mime = EXTENSION_MIME_MAP.get(ext)
        if fallback_mime:
            extractor = EXTRACTORS.get(fallback_mime)

    if not extractor:
        raise ValueError(f"Unsupported file type: {file_type} for file {path}")

    try:
        return extractor(path)
    except Exception as e:
        logger.exception(f"Failed to extract content from {path}: {e}")
        raise


class ChunkSpec(BaseModel):
    """Specification for a single chunk of text."""

    index: int
    text: str
    char_start: int
    char_end: int


def normalize_whitespace(text: str) -> str:
    """Normalize whitespace in text."""
    # Replace multiple whitespace with single space
    text = re.sub(r"[ \t]+", " ", text)
    # Replace multiple newlines with double newline
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(
    text: str,
    chunk_size: int = 1000,
    overlap: int = 200,
) -> list[ChunkSpec]:
    """
    Split text into overlapping chunks.

    Args:
        text: The text to chunk
        chunk_size: Maximum characters per chunk
        overlap: Number of characters to overlap between chunks

    Returns:
        List of ChunkSpec objects
    """
    text = normalize_whitespace(text)

    if not text:
        return []

    chunks = []
    start = 0
    index = 0

    while start < len(text):
        end = min(len(text), start + chunk_size)

        # Try to break at a sentence boundary if not at the end
        if end < len(text):
            # Look for sentence end within the last 20% of the chunk
            search_start = start + int(chunk_size * 0.8)
            for i in range(end, search_start, -1):
                if text[i - 1] in ".!?\n":
                    end = i
                    break

        chunk_text_content = text[start:end]

        chunks.append(
            ChunkSpec(
                index=index,
                text=chunk_text_content,
                char_start=start,
                char_end=end,
            )
        )

        index += 1

        # Move start forward, but ensure we make progress
        new_start = end - overlap
        if new_start <= start:
            new_start = end

        start = new_start

        # Safety check to prevent infinite loops
        if start >= len(text):
            break

    return chunks
